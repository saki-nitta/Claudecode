#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SC講話配布資料（スライド印刷用）生成スクリプト
テーマ：教育相談の在り方 — SCが大切にしていること —
対象：養護実習生（岩手町立川口中学校）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラーパレット ──────────────────────────────────
BLUE    = RGBColor(0x2E, 0x75, 0xB6)
LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)
ORANGE  = RGBColor(0xED, 0x7D, 0x31)
LORANGE = RGBColor(0xFF, 0xF2, 0xCC)
GREEN   = RGBColor(0x37, 0x86, 0x34)
LGREEN  = RGBColor(0xE2, 0xEF, 0xDA)
RED     = RGBColor(0xC0, 0x00, 0x00)
LRED    = RGBColor(0xFC, 0xE4, 0xD6)
DARK    = RGBColor(0x26, 0x26, 0x26)
GRAY    = RGBColor(0x55, 0x55, 0x55)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEAL    = RGBColor(0x17, 0x6A, 0x6A)

# ── ユーティリティ ──────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, border=None, bw=0.75):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, l, t, w, h, text, size=12, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def para(tf, text, size=12, bold=False, italic=False,
         color=DARK, align=PP_ALIGN.LEFT, sb=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(sb)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

def header(slide, title, sub=None, h=1.05):
    rect(slide, 0, 0, 10, h, fill=BLUE)
    if sub:
        txt(slide, 0.25, 0.1,  9.5, 0.55, title, size=20, bold=True, color=WHITE)
        txt(slide, 0.25, 0.65, 9.5, 0.38, sub,   size=11, color=LBLUE)
    else:
        txt(slide, 0.25, 0.22, 9.5, 0.62, title, size=22, bold=True, color=WHITE)

def label(slide, l, t, w, h, text, bg=BLUE, fg=WHITE, size=11, bold=True):
    rect(slide, l, t, w, h, fill=bg)
    txt(slide, l+0.1, t+0.05, w-0.15, h-0.06, text, size=size, bold=bold, color=fg)

def card(slide, l, t, w, h, title, lines, tbg=BLUE, bbg=LGRAY, tsz=12, bsz=11):
    rect(slide, l, t, w, 0.4, fill=tbg)
    txt(slide, l+0.1, t+0.06, w-0.15, 0.3, title, size=tsz, bold=True, color=WHITE)
    rect(slide, l, t+0.4, w, h-0.4, fill=bbg)
    tf = txt(slide, l+0.12, t+0.48, w-0.22, h-0.52, lines[0], size=bsz, color=DARK)
    for line in lines[1:]:
        para(tf, line, size=bsz, color=DARK)

# ── スライド生成 ─────────────────────────────────────

def s01_title(prs):
    """タイトル"""
    s = blank(prs)
    rect(s, 0, 0, 10, 7.5, fill=BLUE)
    rect(s, 0.5, 3.05, 9, 0.06, fill=WHITE)
    txt(s, 0.5, 1.1, 9, 1.5, "教育相談の在り方",
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 3.2, 9, 0.7, "— SCが大切にしていること —",
        size=18, color=LBLUE, align=PP_ALIGN.CENTER)
    txt(s, 1.5, 4.1, 7, 0.7, "気づく・聴く・つなぐ",
        size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 6.55, 9, 0.5,
        "令和8年5月26日　岩手町立川口中学校  スクールカウンセラー",
        size=11, color=LBLUE, align=PP_ALIGN.CENTER)

def s02_goal(prs):
    """今日のゴール"""
    s = blank(prs)
    header(s, "今日のゴール")

    cards_data = [
        ("① 気づく",  "サインを見逃さない\n視点",         BLUE),
        ("② 聴く",    "安心して話せる\n関わり方\n（今日のメイン）", ORANGE),
        ("③ つなぐ",  "SCへの\nつなぎ方",                 GREEN),
    ]
    for i, (lbl, body, c) in enumerate(cards_data):
        l = 0.3 + i * 3.17
        rect(s, l, 1.22, 2.9, 2.7, fill=c)
        txt(s, l+0.1, 1.32, 2.7, 0.65, lbl,
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, l+0.1, 2.05, 2.7, 1.7, body,
            size=13, color=WHITE, align=PP_ALIGN.CENTER)

    rect(s, 0.3, 4.1, 9.4, 0.62, fill=LORANGE)
    txt(s, 0.4, 4.19, 9.2, 0.45,
        "今日は「聴く技術（マイクロカウンセリング）」を中心に扱います",
        size=13, bold=True, color=RGBColor(0x70, 0x30, 0x00),
        align=PP_ALIGN.CENTER)

    rect(s, 0.3, 4.85, 9.4, 0.9, fill=LGRAY)
    tf = txt(s, 0.42, 4.93, 9.1, 0.75,
             "生徒理解・保健室経営・連携については、養護教諭の先生からもお話があります。",
             size=10.5, color=GRAY)
    para(tf, "今日のSC講話では「心理的な聴き方」の視点に特化してお伝えします。",
         size=10.5, color=GRAY)

def s03_hokenshitsu(prs):
    """SCから見た保健室"""
    s = blank(prs)
    header(s, "SCから見た保健室", sub="〜 来室する生徒のサインをどう読むか 〜")

    rect(s, 0.3, 1.28, 4.45, 2.1, fill=LBLUE)
    txt(s, 0.42, 1.35, 4.2, 0.4,
        "身体の訴えの裏にあるもの", size=12, bold=True, color=BLUE)
    tf = txt(s, 0.42, 1.82, 4.2, 1.42, "「頭痛」「腹痛」→ 本当は？",
             size=12, color=DARK)
    para(tf, "・友人関係のストレス",  size=11, color=DARK)
    para(tf, "・家庭内の問題",        size=11, color=DARK)
    para(tf, "・学業プレッシャー",    size=11, color=DARK)

    rect(s, 5.25, 1.28, 4.45, 2.1, fill=LGRAY)
    txt(s, 5.37, 1.35, 4.2, 0.4, "行動のサイン", size=12, bold=True, color=BLUE)
    tf2 = txt(s, 5.37, 1.82, 4.2, 1.42, "・毎日同じ時間に来る", size=11, color=DARK)
    para(tf2, "・話さずに座っている",  size=11)
    para(tf2, "・帰りたがらない",      size=11)
    para(tf2, "・沈黙が多い",          size=11)

    rect(s, 0.3, 3.55, 9.4, 0.65, fill=BLUE)
    txt(s, 0.4, 3.63, 9.2, 0.5,
        "「また来た」→「何かあるかも」　症状ではなく「背景」を見る",
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rect(s, 0.3, 4.35, 9.4, 1.0, fill=LGRAY)
    txt(s, 0.42, 4.43, 2.0, 0.32, "SCの視点：", size=11, bold=True, color=BLUE)
    txt(s, 0.42, 4.78, 9.1, 0.5,
        "その子がこの時間・この場所に来た文脈を、一緒に想像してみる",
        size=12, color=DARK, italic=True)

def s04_kiku(prs):
    """「聴く」とは何か"""
    s = blank(prs)
    header(s, "「聴く」とは何か")

    rect(s, 0.3, 1.18, 4.3, 3.0, fill=LRED)
    txt(s, 0.42, 1.27, 4.1, 0.42, "✗  よくある誤解",
        size=13, bold=True, color=RED)
    tf = txt(s, 0.42, 1.77, 4.1, 2.25, "・アドバイスしなければ", size=12, color=DARK)
    para(tf, "・解決策を教えなければ", size=12)
    para(tf, "・何か言わなければ",     size=12)
    para(tf, "・沈黙はまずい",         size=12)

    txt(s, 4.35, 2.3, 1.3, 0.6, "→", size=28, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)

    rect(s, 5.4, 1.18, 4.3, 3.0, fill=LGREEN)
    txt(s, 5.52, 1.27, 4.1, 0.42, "〇  本当の「聴く」",
        size=13, bold=True, color=GREEN)
    tf2 = txt(s, 5.52, 1.77, 4.1, 2.25, "・話し続けてもらうこと",
              size=12, bold=True, color=DARK)
    para(tf2, "・安心感を届けること", size=12)
    para(tf2, "・判断しないこと",     size=12)
    para(tf2, "・沈黙も関わりのひとつ", size=12)

    rect(s, 0.3, 4.38, 9.4, 0.68, fill=BLUE)
    txt(s, 0.4, 4.47, 9.2, 0.52,
        "「聴く」こと自体が、すでに支援である",
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, 0.3, 5.18, 9.4, 0.55,
        "＊解決できなくても、「聴いてもらえた」という体験が生徒の安心感をつくります。",
        size=11, color=GRAY, italic=True)

def s05_micro(prs):
    """マイクロカウンセリングとは"""
    s = blank(prs)
    header(s, "マイクロカウンセリングとは",
           sub="具体的な「聴く技術」の体系（Ivey, 1971）")

    txt(s, 0.3, 1.28, 9.4, 0.52,
        "カウンセリングで実際に使われる「関わり行動」を体系化したもの。今日は3つの技法を中心に扱います。",
        size=12, color=DARK)

    techs = [
        ("① 最小限のはげまし", "うなずき・「うん」・「そうなんだ」\n話を止めずに促す", BLUE),
        ("② 感情の反映",       "相手の感情をことばにして返す\n「それはつらかったね」",  ORANGE),
        ("③ 開いた質問",       "Yes/Noで終わらない問いかけ\n「どんなことが気になってる？」", GREEN),
    ]
    for i, (title, body, c) in enumerate(techs):
        l = 0.3 + i * 3.17
        rect(s, l, 2.0, 2.9, 0.48, fill=c)
        txt(s, l+0.1, 2.07, 2.7, 0.36, title, size=13, bold=True, color=WHITE)
        rect(s, l, 2.48, 2.9, 1.7, fill=LGRAY)
        txt(s, l+0.12, 2.57, 2.68, 1.55, body, size=11.5, color=DARK)

    rect(s, 0.3, 4.38, 9.4, 0.52, fill=LORANGE)
    txt(s, 0.4, 4.46, 9.2, 0.38,
        "共通のゴール：「話し続けてもらうこと」「安心できる場をつくること」",
        size=12, bold=True, color=RGBColor(0x70, 0x30, 0x00),
        align=PP_ALIGN.CENTER)

    txt(s, 0.3, 5.03, 9.4, 0.7,
        "これらは「カウンセリングの特別な技術」ではなく、今日から保健室でも使える具体的な関わり方です。",
        size=11, color=GRAY)

def s06_minimal(prs):
    """①最小限のはげまし"""
    s = blank(prs)
    header(s, "① 最小限のはげまし",
           sub="Minimal Encouragers ─ 話を止めずに促す最も小さな関わり")

    rect(s, 0.3, 1.3, 4.45, 4.65, fill=LGRAY)
    label(s, 0.3, 1.3, 4.45, 0.42, "言語的なはげまし", bg=BLUE)
    tf = txt(s, 0.5, 1.88, 4.1, 0.6, "「うん」", size=20, bold=True, color=BLUE)
    para(tf, "「そうなんだ」", size=20, bold=True, color=BLUE, sb=10)
    para(tf, "「なるほど」",   size=20, bold=True, color=BLUE, sb=10)
    para(tf, "「それで？」",   size=20, bold=True, color=BLUE, sb=10)
    para(tf, " ",              size=6)
    para(tf, "ポイント：短く・自然に・言いすぎない", size=10, color=GRAY, italic=True)

    rect(s, 5.25, 1.3, 4.45, 4.65, fill=LGRAY)
    label(s, 5.25, 1.3, 4.45, 0.42, "非言語的なはげまし", bg=BLUE)
    tf2 = txt(s, 5.42, 1.88, 4.1, 0.55, "うなずく",
              size=17, bold=True, color=DARK)
    para(tf2, "アイコンタクト（柔らかく）", size=17, bold=True, color=DARK, sb=12)
    para(tf2, "少し前傾姿勢",              size=17, bold=True, color=DARK, sb=12)
    para(tf2, "表情をそろえる",            size=17, bold=True, color=DARK, sb=12)
    para(tf2, " ", size=6)
    para(tf2, "言葉以上に「聴いている」が伝わる", size=10, color=GRAY, italic=True)

    rect(s, 0.3, 6.1, 9.4, 0.52, fill=BLUE)
    txt(s, 0.4, 6.18, 9.2, 0.38,
        "目的：安心して話し続けてもらうこと。解決やアドバイスは不要。",
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def s07_reflection(prs):
    """②感情の反映"""
    s = blank(prs)
    header(s, "② 感情の反映",
           sub="Reflection of Feeling ─ 感情をことばにして返す")

    rect(s, 0.3, 1.3, 9.4, 0.65, fill=LORANGE)
    txt(s, 0.4, 1.38, 9.2, 0.5,
        "相手が感じている感情を、こちらがことばにして伝え返すこと",
        size=14, bold=True, color=RGBColor(0x70, 0x30, 0x00),
        align=PP_ALIGN.CENTER)

    examples = [
        ("友達に無視されてる…",          "それはつらかったんだね",       "悲しみ・孤独"),
        ("テストのことが頭から離れない",  "すごく心配なんだね",           "不安"),
        ("（うつむいたまま、黙っている）","なんか言いにくいことがあるのかな", "緊張・迷い"),
    ]
    for i, (gen, res, emo) in enumerate(examples):
        t = 2.1 + i * 1.28
        rect(s, 0.3,  t, 3.9, 1.1, fill=LGRAY)
        txt(s, 0.42, t+0.1, 3.7, 0.3, "生徒の言葉", size=9, color=GRAY)
        txt(s, 0.42, t+0.42, 3.7, 0.58, f"「{gen}」", size=11, color=DARK)

        txt(s, 4.2, t+0.38, 0.9, 0.42, "→", size=22, bold=True,
            color=ORANGE, align=PP_ALIGN.CENTER)

        rect(s, 5.1, t, 4.6, 1.1, fill=LGREEN)
        txt(s, 5.22, t+0.1, 4.4, 0.3, f"感情：{emo}", size=9, color=GREEN, bold=True)
        txt(s, 5.22, t+0.42, 4.4, 0.58, f"「{res}」", size=12, bold=True, color=DARK)

    rect(s, 0.3, 5.98, 9.4, 0.65, fill=LGRAY)
    txt(s, 0.42, 6.06, 9.1, 0.52,
        "＊感情の名前を「正確に当てる」必要はない。受け止めようとする姿勢が伝わることが大切。",
        size=10.5, color=GRAY, italic=True)

def s08_open(prs):
    """③開いた質問"""
    s = blank(prs)
    header(s, "③ 開いた質問 と 閉じた質問",
           sub="Open / Closed Question ─ 問いの種類で会話の広がり方が変わる")

    rect(s, 0.3, 1.3, 4.4, 0.45, fill=RED)
    txt(s, 0.42, 1.36, 4.2, 0.34, "閉じた質問（Yes/Noで答えられる）",
        size=12, bold=True, color=WHITE)
    rect(s, 5.3, 1.3, 4.4, 0.45, fill=GREEN)
    txt(s, 5.42, 1.36, 4.2, 0.34, "開いた質問（自由に話せる）",
        size=12, bold=True, color=WHITE)

    pairs = [
        ("「友達とけんかした？」",     "「友達のこと、どんなことが気になってる？」"),
        ("「つらい？」",               "「今、どんな気持ちでいる？」"),
        ("「ちゃんと寝てる？」",       "「最近、夜はどんなふうに過ごしてる？」"),
    ]
    for i, (c, o) in enumerate(pairs):
        t = 1.9 + i * 1.08
        rect(s, 0.3, t, 4.4, 0.92, fill=LRED)
        txt(s, 0.42, t+0.15, 4.2, 0.62, c, size=12, color=DARK)
        rect(s, 5.3, t, 4.4, 0.92, fill=LGREEN)
        txt(s, 5.42, t+0.15, 4.2, 0.62, o, size=12, color=DARK)

    rect(s, 0.3, 5.2, 9.4, 0.58, fill=LORANGE)
    txt(s, 0.4, 5.27, 9.2, 0.44,
        "開いた質問は「話す余地」を与える。ただし使いすぎると尋問になるので注意。",
        size=12, bold=True, color=RGBColor(0x70, 0x30, 0x00),
        align=PP_ALIGN.CENTER)
    txt(s, 0.3, 5.9, 9.4, 0.45,
        "バランスが大切：開いた質問で話を広げ、閉じた質問で事実を確認する",
        size=11, color=GRAY)

def s09_roleplay(prs):
    """ミニロールプレイ"""
    s = blank(prs)
    header(s, "ミニロールプレイ", sub="技法を実際に使ってみよう（2分）")

    label(s, 0.3, 1.28, 9.4, 0.42, "場面設定", bg=BLUE)
    rect(s, 0.3, 1.7, 9.4, 1.08, fill=LGRAY)
    tf = txt(s, 0.42, 1.78, 9.1, 0.95,
             "授業中にお腹が痛いといって保健室に来た中2の女子生徒。",
             size=12, color=DARK)
    para(tf, "横になっているうちに少し話せる様子になり、ぽつりとつぶやいた。", size=12)
    para(tf, " ", size=5)
    para(tf, "「…最近、学校来るのがしんどくて」", size=14, bold=True, color=BLUE)

    rect(s, 0.3, 2.9, 9.4, 0.5, fill=ORANGE)
    txt(s, 0.4, 2.98, 9.2, 0.38, "あなたなら、どう返しますか？",
        size=15, bold=True, color=WHITE)

    hints = [
        ("最小限のはげまし", "うなずく\n「そうなんだ」", BLUE),
        ("感情の反映",       "「しんどいんだね」",       ORANGE),
        ("開いた質問",       "「どんなときが\n特にしんどい？」", GREEN),
    ]
    for i, (tech, ex, c) in enumerate(hints):
        l = 0.3 + i * 3.17
        rect(s, l, 3.57, 2.9, 0.42, fill=c)
        txt(s, l+0.1, 3.63, 2.7, 0.3, tech, size=11, bold=True, color=WHITE)
        rect(s, l, 3.99, 2.9, 0.82, fill=LGRAY)
        txt(s, l+0.12, 4.06, 2.68, 0.7, ex, size=11.5, color=DARK)

    rect(s, 0.3, 4.95, 9.4, 0.88, fill=LGRAY)
    txt(s, 0.42, 5.02, 9.1, 0.3, "ロールプレイのポイント",
        size=11, bold=True, color=BLUE)
    tf2 = txt(s, 0.42, 5.33, 9.1, 0.45,
              "・正解はない。「どうしようと思ったか」のプロセスが大事",
              size=11, color=DARK)
    para(tf2, "・やってみた感覚（難しかった／自然だったなど）も大切な気づき", size=11)

def s10_timing(prs):
    """いつSCにつなぐか"""
    s = blank(prs)
    header(s, "いつ・どうやってSCにつなぐか",
           sub="タイミングの見極めと声かけの実際")

    label(s, 0.3, 1.28, 9.4, 0.42, "こんなときはSCへ", bg=BLUE)

    timings = [
        ("継続する悩み",    "同じ問題が何週間も続いている、保健室来室が長期化している"),
        ("深刻なサイン",    "自傷・希死念慮の言及、不登校の始まり、強い情緒不安定"),
        ("家庭環境の問題",  "保護者関係・経済的困難など、学校内だけでは対応が難しい背景"),
        ("本人の希望",      "「もっと話したい」「誰かに聞いてほしい」という意思がある"),
    ]
    for i, (lbl, detail) in enumerate(timings):
        t = 1.85 + i * 0.92
        rect(s, 0.3, t, 2.25, 0.76, fill=LBLUE)
        txt(s, 0.4, t+0.18, 2.1, 0.42, lbl, size=12, bold=True, color=BLUE)
        rect(s, 2.55, t, 7.15, 0.76, fill=LGRAY)
        txt(s, 2.65, t+0.16, 7.0, 0.48, detail, size=11, color=DARK)

    rect(s, 0.3, 5.58, 9.4, 0.4, fill=ORANGE)
    txt(s, 0.4, 5.64, 9.2, 0.3, "SCへつなぐときの声かけ例",
        size=12, bold=True, color=WHITE)
    rect(s, 0.3, 5.98, 9.4, 0.65, fill=LGREEN)
    txt(s, 0.42, 6.06, 9.1, 0.52,
        "「もう少し詳しく話を聞いてくれる先生がいるんだけど、一緒に話してみる？」",
        size=12, bold=True, color=DARK)

def s11_renkei(prs):
    """養護教諭とSCの連携"""
    s = blank(prs)
    header(s, "養護教諭とSCの連携",
           sub="それぞれの強みを活かすチーム支援")

    rect(s, 0.3, 1.28, 4.4, 0.44, fill=TEAL)
    txt(s, 0.42, 1.34, 4.2, 0.34, "養護教諭の強み",
        size=13, bold=True, color=WHITE)
    rect(s, 5.3, 1.28, 4.4, 0.44, fill=BLUE)
    txt(s, 5.42, 1.34, 4.2, 0.34, "SCの強み",
        size=13, bold=True, color=WHITE)

    nurse = ["毎日そこにいる安心感", "身体的ケアとの一体対応",
             "全生徒の状態を把握", "保護者・担任との日常連携"]
    sc    = ["専門的な面談時間（週1〜2回）", "心理アセスメントの視点",
             "守秘義務による安心した開示", "外部機関との連携知識"]

    for i, (n, sc_) in enumerate(zip(nurse, sc)):
        t = 1.86 + i * 0.72
        rect(s, 0.3, t, 4.4, 0.6, fill=LGRAY)
        txt(s, 0.42, t+0.12, 4.2, 0.4, f"・{n}", size=11, color=DARK)
        rect(s, 5.3, t, 4.4, 0.6, fill=LGRAY)
        txt(s, 5.42, t+0.12, 4.2, 0.4, f"・{sc_}", size=11, color=DARK)

    rect(s, 0.3, 4.78, 9.4, 0.62, fill=BLUE)
    txt(s, 0.4, 4.86, 9.2, 0.48,
        "「一人で抱えない」── 情報を共有しながら、チームで支える",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tf = txt(s, 0.3, 5.52, 9.4, 0.42,
             "情報共有のポイント：「事実」と「SCの解釈」を分けて伝える",
             size=10.5, color=GRAY)
    para(tf, "守秘義務：原則として本人の同意のもとで共有（緊急性がある場合は除く）",
         size=10.5, color=GRAY)

def s12_matome(prs):
    """まとめ"""
    s = blank(prs)
    header(s, "今日のまとめ")

    techs = [
        ("① 最小限のはげまし", "うなずき・「うん」「そうなんだ」で話を促す",      BLUE),
        ("② 感情の反映",       "「それはつらかったね」感情をことばにして返す",    ORANGE),
        ("③ 開いた質問",       "「どんなことが気になってる？」話す余地をつくる",  GREEN),
    ]
    for i, (title, body, c) in enumerate(techs):
        t = 1.28 + i * 1.38
        rect(s, 0.3, t, 9.4, 1.22, fill=LGRAY)
        rect(s, 0.3, t, 2.7, 1.22, fill=c)
        txt(s, 0.4, t+0.35, 2.52, 0.55, title,
            size=14, bold=True, color=WHITE)
        txt(s, 3.15, t+0.33, 6.3, 0.6, body, size=13, color=DARK)

    rect(s, 0.3, 4.42, 9.4, 0.6, fill=BLUE)
    txt(s, 0.4, 4.5, 9.2, 0.45,
        "共通のゴール：「安心して話し続けてもらうこと」",
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, 0.3, 5.15, 9.4, 0.55,
        "この3つは「特別なカウンセリング技術」ではなく、今日から保健室でも使える関わり方です。",
        size=12, color=GRAY)

def s13_closing(prs):
    """最後のことば"""
    s = blank(prs)
    rect(s, 0, 0, 10, 7.5, fill=BLUE)
    rect(s, 0.5, 2.95, 9, 0.06, fill=WHITE)
    rect(s, 0.5, 5.2, 9, 0.06, fill=WHITE)

    txt(s, 0.5, 0.9, 9, 1.8,
        "保健室は\n「問題を解決する場所」ではなく、",
        size=22, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 3.1, 9, 1.8,
        "「安心して出発できる場所」",
        size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 5.35, 9, 0.8,
        "その場所を一緒につくっていきましょう。",
        size=18, color=LBLUE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 6.55, 9, 0.5,
        "川口中学校  スクールカウンセラー",
        size=12, color=LBLUE, align=PP_ALIGN.CENTER)

# ── メイン ──────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    s01_title(prs)
    s02_goal(prs)
    s03_hokenshitsu(prs)
    s04_kiku(prs)
    s05_micro(prs)
    s06_minimal(prs)
    s07_reflection(prs)
    s08_open(prs)
    s09_roleplay(prs)
    s10_timing(prs)
    s11_renkei(prs)
    s12_matome(prs)
    s13_closing(prs)

    out = r"d:\toumei.shippo\20260513_doc_SC講話教育相談_draft_v1.pptx"
    prs.save(out)
    print(f"完了：{out}")
    print(f"スライド数：{len(prs.slides)}")

if __name__ == "__main__":
    main()
