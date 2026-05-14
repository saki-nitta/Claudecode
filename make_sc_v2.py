#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""SC講話配布資料 v2 - 白黒印刷対応デザイン"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE   = RGBColor(0x2E,0x75,0xB6); LBLUE  = RGBColor(0xBD,0xD7,0xEE)
ORANGE = RGBColor(0xED,0x7D,0x31); GREEN  = RGBColor(0x37,0x86,0x34)
DARK   = RGBColor(0x26,0x26,0x26); GRAY   = RGBColor(0x55,0x55,0x55)
BGRAY  = RGBColor(0xBB,0xBB,0xBB); WHITE  = RGBColor(0xFF,0xFF,0xFF)
TEAL   = RGBColor(0x17,0x6A,0x6A); DRED   = RGBColor(0x99,0x00,0x00)
DGREEN = RGBColor(0x20,0x70,0x20); CRED   = RGBColor(0xC0,0x00,0x00)

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def frect(s, l, t, w, h, fill, bdr=None, bw=0.5):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if bdr: sh.line.color.rgb = bdr; sh.line.width = Pt(bw)
    else: sh.line.fill.background()

def orect(s, l, t, w, h, bdr=None, bw=0.75):
    c = bdr if bdr else BGRAY
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.background(); sh.line.color.rgb = c; sh.line.width = Pt(bw)

def lbar(s, l, t, h, c=None):
    frect(s, l, t, 0.07, h, fill=(c if c else BLUE))

def tx(s, l, t, w, h, text, sz=12, bd=False, it=False, col=None, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text; r.font.size = Pt(sz)
    r.font.bold = bd; r.font.italic = it
    r.font.color.rgb = col if col else DARK
    return tf

def pr(tf, text, sz=12, bd=False, it=False, col=None, sb=4):
    p = tf.add_paragraph(); p.space_before = Pt(sb)
    r = p.add_run(); r.text = text; r.font.size = Pt(sz)
    r.font.bold = bd; r.font.italic = it
    r.font.color.rgb = col if col else DARK

def hl(s, l, t, w):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.01))
    sh.fill.solid(); sh.fill.fore_color.rgb = BGRAY; sh.line.fill.background()

def hd(s, title, sub=None, h=1.05):
    frect(s, 0, 0, 10, h, fill=BLUE)
    if sub:
        tx(s, 0.25, 0.1,  9.5, 0.55, title, sz=20, bd=True, col=WHITE)
        tx(s, 0.25, 0.65, 9.5, 0.38, sub,   sz=11, col=LBLUE)
    else:
        tx(s, 0.25, 0.22, 9.5, 0.62, title, sz=22, bd=True, col=WHITE)

def slbl(s, l, t, w, h, text, bg=None, sz=11):
    frect(s, l, t, w, h, fill=(bg if bg else BLUE))
    tx(s, l+0.1, t+0.05, w-0.15, h-0.06, text, sz=sz, bd=True, col=WHITE)

def s01(prs):
    s = blank(prs)
    frect(s, 0, 0, 10, 3.5, fill=BLUE)
    tx(s, 0.5, 0.6,  9, 1.6, "教育相談の在り方",
       sz=38, bd=True, col=WHITE, al=PP_ALIGN.CENTER)
    tx(s, 0.5, 2.2,  9, 0.7, "SCが大切にしていること",
       sz=18, col=LBLUE, al=PP_ALIGN.CENTER)
    tx(s, 1.5, 3.9,  7, 0.8,
       "気づく・聴く・つなぐ",
       sz=26, bd=True, col=BLUE, al=PP_ALIGN.CENTER)
    hl(s, 1.5, 4.82, 7)
    tx(s, 0.5, 5.1,  9, 0.55,
       "令和8年5月26日　岩手町立川口中学校　スクールカウンセラー",
       sz=12, col=GRAY, al=PP_ALIGN.CENTER)

def s02(prs):
    s = blank(prs); hd(s, "今日のゴール")
    labels = [
        ("① 気づく",
         "サインを見逃さない視点", BLUE),
        ("② 聴く",
         "安心して話せる関わり方\n（今日のメイン）", ORANGE),
        ("③ つなぐ",
         "SCへのつなぎ方", GREEN),
    ]
    for i, (lbl, body, c) in enumerate(labels):
        l = 0.3 + i*3.17
        slbl(s, l, 1.22, 2.9, 0.48, lbl, bg=c, sz=14)
        orect(s, l, 1.7,  2.9, 1.6)
        tx(s, l+0.14, 1.83, 2.65, 1.35, body, sz=12)
    lbar(s, 0.3, 3.5, 0.62, c=ORANGE)
    tx(s, 0.48, 3.55, 8.9, 0.52,
       "今日は「聴く技術（マイクロカウンセリング）」を中心に扱います",
       sz=13, bd=True)
    hl(s, 0.3, 4.28, 9.4)
    tf = tx(s, 0.3, 4.38, 9.4, 0.42,
            "生徒理解・保健室経営・連携については、養護教諭の先生からもお話があります。",
            sz=10.5, col=GRAY)
    pr(tf,
       "今日のSC講話では「心理的な聴き方」の視点に特化してお伝えします。",
       sz=10.5, col=GRAY)

def s03(prs):
    s = blank(prs)
    hd(s, "SCから見た保健室",
       sub="来室する生徒のサインをどう読むか")
    slbl(s, 0.3,  1.28, 4.45, 0.38,
         "身体の訴えの裏にあるもの", sz=11)
    orect(s, 0.3,  1.66, 4.45, 1.65)
    tf = tx(s, 0.44, 1.78, 4.2, 1.45,
            "「頭痛」「肅痛」→ 本当は？",
            sz=12, bd=True)
    pr(tf, "・友人関係のストレス", sz=11)
    pr(tf, "・家庭内の問題",                   sz=11)
    pr(tf, "・学業プレッシャー",       sz=11)
    slbl(s, 5.25, 1.28, 4.45, 0.38, "行動のサイン", sz=11)
    orect(s, 5.25, 1.66, 4.45, 1.65)
    tf2 = tx(s, 5.4, 1.78, 4.2, 1.45,
             "・毎日同じ時間帯に来る", sz=11)
    pr(tf2, "・話さずに座っている",   sz=11)
    pr(tf2, "・帰りたがらない",               sz=11)
    pr(tf2, "・沈黙が多い",                           sz=11)
    lbar(s, 0.3, 3.5, 0.62)
    tx(s, 0.48, 3.55, 9.1, 0.52,
       "「また来た」ではなく →「何かあるかも」  症状ではなく「背景」を見る",
       sz=14, bd=True)
    hl(s, 0.3, 4.24, 9.4)
    tx(s, 0.3,  4.34, 1.5, 0.35,
       "SCの視点：", sz=11, bd=True, col=BLUE)
    tx(s, 1.85, 4.34, 7.8, 0.6,
       "その子がこの時間・この場所に来た文脈を、一緒に想像してみる",
       sz=12, it=True)

def s04(prs):
    s = blank(prs)
    hd(s, "「聴く」とは何か")
    slbl(s, 0.3, 1.18, 4.3, 0.38,
         "✗  よくある誤解", bg=DRED, sz=12)
    orect(s, 0.3, 1.56, 4.3, 2.65, bdr=CRED, bw=1.5)
    tf = tx(s, 0.44, 1.7, 4.0, 2.4,
            "・アドバイスしなければ", sz=12)
    pr(tf, "・解決策を教えなければ", sz=12)
    pr(tf, "・何か言わなければ",             sz=12)
    pr(tf, "・沈黙はまずい",                         sz=12)
    tx(s, 4.35, 2.55, 1.3, 0.55, "→",
       sz=28, bd=True, col=ORANGE, al=PP_ALIGN.CENTER)
    slbl(s, 5.4, 1.18, 4.3, 0.38,
         "○  本当の「聴く」", bg=DGREEN, sz=12)
    orect(s, 5.4, 1.56, 4.3, 2.65, bdr=GREEN, bw=1.5)
    tf2 = tx(s, 5.54, 1.7, 4.0, 2.4,
             "・話し続けてもらうこと",
             sz=12, bd=True)
    pr(tf2, "・安心感を届けること",     sz=12)
    pr(tf2, "・判断しないこと",                  sz=12)
    pr(tf2, "・沈黙も関わりのひとつ", sz=12)
    lbar(s, 0.3, 4.4, 0.62)
    tx(s, 0.48, 4.47, 9.1, 0.52,
       "「聴く」こと自体が、すでに支援である",
       sz=16, bd=True)
    tx(s, 0.3, 5.15, 9.4, 0.55,
       "＊解決できなくても、「聴いてもらえた」という体験が生徒の安心感をつくります。",
       sz=11, it=True, col=GRAY)

def s05(prs):
    s = blank(prs)
    hd(s, "マイクロカウンセリングとは",
       sub="具体的な「聴く技術」の体系（Ivey, 1971）")
    tx(s, 0.3, 1.28, 9.4, 0.48,
       "カウンセリングで実際に使われる「関わり行動」を体系化したもの。今日は3つの技法を中心に扱います。",
       sz=12)
    hl(s, 0.3, 1.82, 9.4)
    techs = [
        ("① 最小限のはげまし",
         "うなずき・「うん」・「そうなんだ」\n話を止めずに促す", BLUE),
        ("② 感情の反映",
         "相手の感情をことばにして返す\n「それはつらかったね」", ORANGE),
        ("③ 開いた質問",
         "Yes/Noで終わらない問いかけ\n「どんなことが気になってる？」", GREEN),
    ]
    for i, (title, body, c) in enumerate(techs):
        l = 0.3 + i*3.17
        slbl(s, l, 1.95, 2.9, 0.42, title, bg=c, sz=12)
        orect(s, l, 2.37, 2.9, 1.8)
        tx(s, l+0.12, 2.5, 2.68, 1.6, body, sz=11.5)
    hl(s, 0.3, 4.32, 9.4)
    lbar(s, 0.3, 4.43, 0.52, c=ORANGE)
    tx(s, 0.48, 4.48, 9.1, 0.42,
       "共通のゴール：「話し続けてもらうこと」「安心できる場をつくること」",
       sz=12, bd=True)
    tx(s, 0.3, 5.05, 9.4, 0.6,
       "これらは「カウンセリングの特別な技術」ではなく、今日から保健室でも使える具体的な関わり方です。",
       sz=11, col=GRAY)

def s06(prs):
    s = blank(prs)
    hd(s, "① 最小限のはげまし",
       sub="Minimal Encouragers  話を止めずに促す最も小さな関わり")
    slbl(s, 0.3,  1.28, 4.45, 0.38,
         "言語的なはげまし", sz=12)
    orect(s, 0.3,  1.66, 4.45, 4.28)
    tf = tx(s, 0.5, 1.82, 4.1, 0.6,
            "「うん」", sz=20, bd=True, col=BLUE)
    pr(tf, "「そうなんだ」",
       sz=20, bd=True, col=BLUE, sb=10)
    pr(tf, "「なるほど」",
       sz=20, bd=True, col=BLUE, sb=10)
    pr(tf, "「それで？」",
       sz=20, bd=True, col=BLUE, sb=10)
    pr(tf, " ", sz=6)
    pr(tf, "ポイント：短く・自然に・言いすぎない",
       sz=10, it=True, col=GRAY)
    slbl(s, 5.25, 1.28, 4.45, 0.38,
         "非言語的なはげまし", sz=12)
    orect(s, 5.25, 1.66, 4.45, 4.28)
    tf2 = tx(s, 5.42, 1.82, 4.1, 0.6,
             "うなずく", sz=17, bd=True)
    pr(tf2, "アイコンタクト（柔らかく）",
       sz=17, bd=True, sb=12)
    pr(tf2, "少し前傾姿勢",  sz=17, bd=True, sb=12)
    pr(tf2, "表情をそろえる", sz=17, bd=True, sb=12)
    pr(tf2, " ", sz=6)
    pr(tf2,
       "言葉以上に「聴いている」が伝わる",
       sz=10, it=True, col=GRAY)
    lbar(s, 0.3, 6.1, 0.52)
    tx(s, 0.48, 6.17, 9.1, 0.38,
       "目的：安心して話し続けてもらうこと。解決やアドバイスは不要。",
       sz=12, bd=True)

def s07(prs):
    s = blank(prs)
    hd(s, "② 感情の反映",
       sub="Reflection of Feeling  感情をことばにして返す")
    lbar(s, 0.3, 1.28, 0.62)
    tx(s, 0.48, 1.33, 9.1, 0.52,
       "相手が感じている感情を、こちらがことばにして伝え返すこと",
       sz=14, bd=True)
    hl(s, 0.3, 1.98, 9.4)
    examples = [
        ("友達に無視されてる…",
         "それはつらかったんだね",
         "悔しみ・孤独"),
        ("テストのことが頭から離れない",
         "すごく心配なんだね",
         "不安"),
        ("（うつむいたまま、黙っている）",
         "なんか言いにくいことがあるのかな",
         "緊張・迷い"),
    ]
    for i, (gen, res, emo) in enumerate(examples):
        t = 2.12 + i*1.25
        orect(s, 0.3, t, 3.9, 1.08)
        tx(s, 0.42, t+0.08, 3.7, 0.28,
           "生徒の言葉", sz=9, col=GRAY)
        tx(s, 0.42, t+0.4,  3.7, 0.58,
           "「" + gen + "」", sz=11)
        tx(s, 4.2, t+0.36, 0.9, 0.42, "→",
           sz=22, bd=True, col=ORANGE, al=PP_ALIGN.CENTER)
        slbl(s, 5.1, t, 4.6, 0.38,
             "感情：" + emo, bg=GREEN, sz=9)
        orect(s, 5.1, t+0.38, 4.6, 0.7)
        tx(s, 5.22, t+0.46, 4.4, 0.52,
           "「" + res + "」", sz=12, bd=True)
    hl(s, 0.3, 5.9, 9.4)
    tx(s, 0.3, 5.98, 9.4, 0.62,
       "＊感情の名前を「正確に当てる」必要はない。受け止めようとする姿勢が伝わることが大切。",
       sz=10.5, it=True, col=GRAY)

def s08(prs):
    s = blank(prs)
    hd(s, "③ 開いた質問 と 閉じた質問",
       sub="Open / Closed Question  問いの種類で会話の広がりが変わる")
    slbl(s, 0.3, 1.28, 4.4, 0.4,
         "閉じた質問（Yes/Noで答えられる）",
         bg=DRED, sz=11)
    slbl(s, 5.3, 1.28, 4.4, 0.4,
         "開いた質問（自由に話せる）",
         bg=DGREEN, sz=11)
    pairs = [
        ("「友達とけんかした？」",
         "「友達のこと、どんなことが気になってる？」"),
        ("「つらい？」",
         "「今、どんな気持ちでいる？」"),
        ("「ちゃんと寝てる？」",
         "「最近、夜はどんなふうに過ごしてる？」"),
    ]
    for i, (c, o) in enumerate(pairs):
        t = 1.84 + i*1.1
        orect(s, 0.3, t, 4.4, 0.95, bdr=CRED,  bw=1.0)
        tx(s, 0.42, t+0.18, 4.2, 0.62, c, sz=12)
        orect(s, 5.3, t, 4.4, 0.95, bdr=GREEN, bw=1.0)
        tx(s, 5.42, t+0.18, 4.2, 0.62, o, sz=12)
    hl(s, 0.3, 5.22, 9.4)
    lbar(s, 0.3, 5.35, 0.52, c=ORANGE)
    tx(s, 0.48, 5.4,  9.1, 0.44,
       "開いた質問は「話す余地」を与える。ただし使いすぎると尋問になるので注意。",
       sz=12, bd=True)
    tx(s, 0.3, 5.98, 9.4, 0.42,
       "バランスが大切：開いた質問で話を広げ、閉じた質問で事実を確認する",
       sz=11, col=GRAY)

def s09(prs):
    s = blank(prs)
    hd(s, "ミニロールプレイ",
       sub="技法を実際に使ってみよう（2分）")
    slbl(s, 0.3, 1.28, 9.4, 0.38,
         "場面設定", sz=12)
    orect(s, 0.3, 1.66, 9.4, 1.08)
    tf = tx(s, 0.44, 1.74, 9.1, 0.9,
            "授業中にお腹が痛いといって保健室に来た中2の女子生徒。",
            sz=12)
    pr(tf,
       "横になっているうちに少し話せる様子になり、ぼつりとつぶやいた。",
       sz=12)
    pr(tf,
       "「…最近、学校来るのがしんどくて」",
       sz=13, bd=True, col=BLUE)
    lbar(s, 0.3, 2.87, 0.55, c=ORANGE)
    tx(s, 0.48, 2.92, 9.1, 0.42,
       "あなたなら、どう返しますか？　→　書いてみましょう",
       sz=14, bd=True)
    hints = [
        ("① 最小限のはげまし", BLUE),
        ("② 感情の反映",                   ORANGE),
        ("③ 開いた質問",                   GREEN),
    ]
    for i, (tech, c) in enumerate(hints):
        l = 0.3 + i*3.17
        slbl(s, l, 3.55, 2.9, 0.38, tech, bg=c, sz=11)
        orect(s, l, 3.93, 2.9, 1.0, bdr=BGRAY, bw=1.0)
    hl(s, 0.3, 5.05, 9.4)
    tx(s, 0.3, 5.13, 9.4, 0.3,
       "ロールプレイのポイント：",
       sz=11, bd=True, col=BLUE)
    tf3 = tx(s, 0.3, 5.46, 9.4, 0.42,
             "・正解はない。「どうしようと思ったか」のプロセスが大事",
             sz=11)
    pr(tf3,
       "・やってみた感覚（難しかった／自然だったなど）も大切な気づき",
       sz=11)

def s10(prs):
    s = blank(prs)
    hd(s, "いつ・どうやってSCにつなぐか",
       sub="タイミングの見極めと声かけの実際")
    slbl(s, 0.3, 1.28, 9.4, 0.38,
         "こんなときはSCへ")
    timings = [
        ("継続する悩み",
         "同じ問題が何週間も続いている、保健室来室が長期化している"),
        ("深刻なサイン",
         "自傷・希死念慮の言及、不登校の始まり、強い情緒不安定"),
        ("家庭環境の問題",
         "保護者関係・経済的困難など → SCまたはSSW（スクールソーシャルワーカー）に相談"),
        ("本人の希望",
         "「もっと話したい」「誰かに聆いてほしい」という意思がある"),
    ]
    for i, (lbl, detail) in enumerate(timings):
        t = 1.82 + i*0.93
        orect(s, 0.3, t, 2.2, 0.78)
        tx(s, 0.4, t+0.19, 2.05, 0.44, lbl, sz=12, bd=True, col=BLUE)
        orect(s, 2.5, t, 7.2, 0.78)
        tx(s, 2.62, t+0.17, 6.95, 0.52, detail, sz=11)
    lbar(s, 0.3, 5.62, 0.65, c=ORANGE)
    tx(s, 0.48, 5.67, 9.1, 0.28,
       "SSW（スクールソーシャルワーカー）とは？",
       sz=11, bd=True)
    tf2 = tx(s, 0.48, 5.98, 9.1, 0.52,
             "福祉の視点から家庭・地域・関係機関と学校をつなぐ専門職。",
             sz=11)
    pr(tf2,
       "経済・虹待・養育など、環境面からのアプローチが必要なときに連携します。",
       sz=11)
    hl(s, 0.3, 6.6, 9.4)
    tx(s, 0.3, 6.67, 9.4, 0.45,
       "声かけ例：「もう少し詳しく聆いてくれる先生がいるんだけど、一緒に話してみる？」",
       sz=11, it=True, col=GRAY)

def s11(prs):
    s = blank(prs)
    hd(s, "養護教諭とSCの連携",
       sub="それぞれの強みを活かすチーム支援")
    slbl(s, 0.3, 1.28, 4.4, 0.4,
         "養護教諭の強み", bg=TEAL)
    slbl(s, 5.3, 1.28, 4.4, 0.4, "SCの強み", bg=BLUE)
    nurse = ["毎日そこにいる安心感",
             "身体的ケアとの一体対応",
             "全生徒の状態を把握",
             "保護者・担任との日常連携"]
    sc = ["専門的な面談時間（週１～２回）",
          "心理アセスメントの視点",
          "守秘義務による安心した開示",
          "外部機関との連携知識"]
    for i, (n, sc_) in enumerate(zip(nurse, sc)):
        t = 1.84 + i*0.72
        orect(s, 0.3, t, 4.4, 0.62)
        tx(s, 0.42, t+0.13, 4.2, 0.42, "・" + n, sz=11)
        orect(s, 5.3, t, 4.4, 0.62)
        tx(s, 5.42, t+0.13, 4.2, 0.42, "・" + sc_, sz=11)
    hl(s, 0.3, 4.78, 9.4)
    lbar(s, 0.3, 4.9, 0.62)
    tx(s, 0.48, 4.97, 9.1, 0.48,
       "「一人で抱えない」  情報を共有しながら、チームで支える",
       sz=14, bd=True)
    tf = tx(s, 0.3, 5.62, 9.4, 0.38,
            "情報共有のポイント：「事実」と「SCの解釈」を分けて伝える",
            sz=10.5, col=GRAY)
    pr(tf,
       "守秘義務：原則として本人の同意のもとで共有（緊急性がある場合は除く）",
       sz=10.5, col=GRAY)

def s12(prs):
    s = blank(prs)
    hd(s, "今日のまとめ")
    techs = [
        ("① 最小限のはげまし",
         "うなずき・「うん」「そうなんだ」で話を促す",       BLUE),
        ("② 感情の反映",
         "「それはつらかったね」感情をことばにして返す", ORANGE),
        ("③ 開いた質問",
         "「どんなことが気になってる？」話す余地をつくる", GREEN),
    ]
    for i, (title, body, c) in enumerate(techs):
        t = 1.28 + i*1.4
        slbl(s, 0.3, t, 2.7, 1.24, title, bg=c, sz=14)
        orect(s, 3.0, t, 6.7, 1.24)
        tx(s, 3.14, t+0.38, 6.42, 0.55, body, sz=13)
    hl(s, 0.3, 5.5, 9.4)
    lbar(s, 0.3, 5.62, 0.62)
    tx(s, 0.48, 5.68, 9.1, 0.48,
       "共通のゴール：「安心して話し続けてもらうこと」",
       sz=15, bd=True)
    tx(s, 0.3, 6.3, 9.4, 0.45,
       "この3つは今日から保健室でも使える関わり方です。",
       sz=12, col=GRAY)

def s13(prs):
    s = blank(prs)
    frect(s, 0, 0, 10, 1.6, fill=BLUE)
    hl(s, 0.5, 1.75, 9)
    tx(s, 0.5, 2.0, 9, 1.2,
       "保健室は\n「問題を解決する場所」ではなく、",
       sz=22, al=PP_ALIGN.CENTER)
    tx(s, 0.5, 3.3, 9, 1.0,
       "「安心して出発できる場所」",
       sz=28, bd=True, col=BLUE, al=PP_ALIGN.CENTER)
    hl(s, 0.5, 4.45, 9)
    tx(s, 0.5, 4.65, 9, 0.75,
       "その場所を一緒につくっていきましょう。",
       sz=18, col=GRAY, al=PP_ALIGN.CENTER)
    tx(s, 0.5, 6.5, 9, 0.5,
       "川口中学校　スクールカウンセラー",
       sz=12, col=GRAY, al=PP_ALIGN.CENTER)

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    for fn in [s01,s02,s03,s04,s05,s06,s07,s08,s09,s10,s11,s12,s13]:
        fn(prs)
    out = r"d:\toumei.shippo\20260513_doc_SC講話教育相談_draft_v2.pptx"
    prs.save(out)
    print("done: " + out)
    print("slides: " + str(len(prs.slides)))

if __name__ == "__main__":
    main()